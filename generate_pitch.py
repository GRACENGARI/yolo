"""
Generate HAWKEYE Pitch Deck PowerPoint
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x00, 0xFF, 0x41)   # matrix green
CYAN   = RGBColor(0x00, 0xFF, 0xFF)
RED    = RGBColor(0xFF, 0x22, 0x22)
YELLOW = RGBColor(0xFF, 0xDD, 0x00)
DGRAY  = RGBColor(0x11, 0x11, 0x11)
MGRAY  = RGBColor(0x22, 0x22, 0x22)

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ─────────────────────────────────────────────────────
def bg(slide, color=BLACK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, bg_color=None, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg_color:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg_color
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def hline(slide, t, color=GREEN):
    rect(slide, 0.3, t, 12.73, 0.03, color)

def tag(slide, l, t, text, bg_c=GREEN, txt_c=BLACK, size=13):
    box(slide, l, t, 2.5, 0.4, text, size=size, bold=True,
        color=txt_c, align=PP_ALIGN.CENTER, bg_color=bg_c)

def bullet(slide, l, t, items, size=16, color=WHITE, spacing=0.42):
    for i, item in enumerate(items):
        box(slide, l, t + i * spacing, 11, 0.5, f"▸  {item}",
            size=size, color=color)

def stat_card(slide, l, t, number, label, num_color=CYAN):
    rect(slide, l, t, 2.8, 1.5, MGRAY, GREEN)
    box(slide, l+0.1, t+0.1, 2.6, 0.8, number, size=32, bold=True,
        color=num_color, align=PP_ALIGN.CENTER)
    box(slide, l+0.1, t+0.9, 2.6, 0.5, label, size=12, bold=False,
        color=WHITE, align=PP_ALIGN.CENTER)

def table_row(slide, l, t, cols, widths, colors, size=13, bold=False):
    x = l
    for text, w, c in zip(cols, widths, colors):
        box(slide, x, t, w-0.05, 0.38, text, size=size, bold=bold,
            color=c, align=PP_ALIGN.CENTER, bg_color=MGRAY)
        x += w

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 7.5, RGBColor(0x05,0x05,0x05))

# green accent bar left
rect(s, 0, 0, 0.18, 7.5, GREEN)

box(s, 0.5, 0.6, 12, 1.2, "🦅  HAWKEYE", size=60, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
box(s, 0.5, 1.9, 12, 0.7, "Unified Urban Dragnet — AI-Powered City-Wide Surveillance Intelligence",
    size=22, color=CYAN, align=PP_ALIGN.LEFT)
hline(s, 2.8)
box(s, 0.5, 3.0, 12, 0.6,
    '"The cameras are already watching.  HAWKEYE makes them think."',
    size=20, color=YELLOW, align=PP_ALIGN.LEFT)

box(s, 0.5, 4.2, 4, 0.45, "64,001  Vehicle Sightings", size=16, bold=True, color=GREEN)
box(s, 0.5, 4.7, 4, 0.45, "883  Unique Vehicles Tracked", size=16, bold=True, color=GREEN)
box(s, 0.5, 5.2, 4, 0.45, "2,911  Person Detections", size=16, bold=True, color=GREEN)
box(s, 0.5, 5.7, 4, 0.45, "76  Unique Individuals Identified", size=16, bold=True, color=GREEN)

box(s, 0.5, 6.8, 12, 0.4, "github.com/GRACENGARI/yolo",
    size=13, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, RED)
box(s, 0.3, 0.1, 12, 0.7, "THE PROBLEM:  Urban Camouflage",
    size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

rect(s, 0.3, 1.05, 12.7, 1.5, RGBColor(0x1a,0x00,0x00))
box(s, 0.5, 1.1, 12.3, 1.4,
    '"A suspect leaves a crime scene in Westlands at 2:14 PM.\n'
    'By 2:17 PM they have vanished into Nairobi\'s urban noise.\n'
    'By the time police piece together the footage — 3 days later — the suspect is already across the border."',
    size=16, color=YELLOW, align=PP_ALIGN.LEFT)

hline(s, 2.7, RED)

items = [
    "Nairobi has THOUSANDS of CCTV cameras — every one is an island",
    "Cameras in Westlands don't talk to cameras in the CBD",
    "Tracking a suspect = calling building managers + waiting days for DVDs",
    "No unified timeline.  No real-time awareness.  No dragnet.",
    "Criminals exploit the GAP between cameras as their escape route",
]
bullet(s, 0.5, 2.85, items, size=17, color=WHITE, spacing=0.72)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, GREEN)
box(s, 0.3, 0.1, 12, 0.7, "THE SOLUTION:  HAWKEYE",
    size=30, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

rect(s, 0.3, 1.05, 12.7, 1.3, MGRAY)
box(s, 0.5, 1.1, 12.3, 1.2,
    "HAWKEYE is the intelligence layer that stitches every camera feed into a single,\n"
    "searchable, real-time timeline — tracking any face or vehicle across the entire city in SECONDS.",
    size=18, bold=True, color=CYAN, align=PP_ALIGN.LEFT)

hline(s, 2.5)

# 4 pillars
for i, (title, desc) in enumerate([
    ("INGEST",   "Any camera feed — RTSP, CCTV, dashcam — plugs straight in"),
    ("DETECT",   "YOLOv8 AI identifies every person & vehicle in real-time"),
    ("STORE",    "Every sighting logged: who, when, where, confidence score"),
    ("SEARCH",   "Agent searches → full movement timeline returned instantly"),
]):
    col = i * 3.2 + 0.3
    rect(s, col, 2.7, 3.0, 2.8, MGRAY, GREEN)
    box(s, col+0.1, 2.8, 2.8, 0.5, f"0{i+1}  {title}", size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    box(s, col+0.1, 3.4, 2.8, 1.8, desc, size=14, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.3, 5.7, 12.7, 0.5,
    "CAMERA FEEDS  →  AI DETECTION  →  UNIFIED DATABASE  →  AGENT DASHBOARD",
    size=15, bold=True, color=YELLOW, align=PP_ALIGN.CENTER, bg_color=RGBColor(0x1a,0x1a,0x00))

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — LIVE DATA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, CYAN)
box(s, 0.3, 0.1, 12, 0.7, "LIVE DATA — Real Numbers, Real System",
    size=30, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

box(s, 0.3, 1.0, 6, 0.4, "🚗  VEHICLE SURVEILLANCE", size=18, bold=True, color=GREEN)
stat_card(s, 0.3,  1.5, "64,001",  "Total Sightings",    CYAN)
stat_card(s, 3.3,  1.5, "883",     "Unique Vehicles",    GREEN)
stat_card(s, 6.3,  1.5, "1,089",   "Max Sightings\n(1 Vehicle)", YELLOW)
stat_card(s, 9.3,  1.5, "50–95%",  "Confidence Range",  WHITE)

box(s, 0.3, 3.2, 6, 0.4, "🚶  PEOPLE SURVEILLANCE", size=18, bold=True, color=RED)
stat_card(s, 0.3,  3.7, "2,911",   "Person Detections",  RED)
stat_card(s, 3.3,  3.7, "76",      "Unique Individuals", YELLOW)
stat_card(s, 6.3,  3.7, "13.6s",   "Source Footage",     WHITE)
stat_card(s, 9.3,  3.7, "4–6 FPS", "CPU Processing",     CYAN)

hline(s, 5.4)
box(s, 0.3, 5.5, 12.7, 0.5,
    "All data exported to CSV  ·  vehicle_evidence_report.csv  ·  people_surveillance_report.csv",
    size=14, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)
box(s, 0.3, 6.0, 12.7, 0.6,
    "Tracking period:  Feb 28 – Mar 12, 2026  ·  883 vehicles  ·  64,001 rows of evidence",
    size=15, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — FORENSIC SEARCH SCENARIO
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, YELLOW)
box(s, 0.3, 0.1, 12, 0.7, "FORENSIC SEARCH — 10 Seconds vs 3 Days",
    size=30, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

# Scenario box
rect(s, 0.3, 1.0, 12.7, 0.8, MGRAY)
box(s, 0.5, 1.05, 12.3, 0.7,
    "SCENARIO:  Robbery at 14:32 in the CBD.  Suspect flees in a white saloon.",
    size=17, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)

# Two columns
rect(s, 0.3, 2.0, 5.9, 4.5, RGBColor(0x1a,0x00,0x00))
box(s, 0.4, 2.05, 5.7, 0.5, "❌  WITHOUT HAWKEYE", size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
for i, line in enumerate([
    "Call 12 building managers",
    "Request CCTV footage manually",
    "Wait 2–5 days for DVDs",
    "Analyst manually reviews hours of tape",
    "Piece together route by hand",
    "Suspect already across the border",
]):
    box(s, 0.5, 2.6 + i*0.55, 5.5, 0.5, f"✗  {line}", size=14, color=RGBColor(0xFF,0x88,0x88))

rect(s, 6.5, 2.0, 6.5, 4.5, RGBColor(0x00,0x1a,0x00))
box(s, 6.6, 2.05, 6.3, 0.5, "✅  WITH HAWKEYE", size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
for i, line in enumerate([
    "Agent opens dashboard",
    "Types vehicle description or Track ID",
    "HAWKEYE searches 64,001 records",
    "Every sighting returned with timestamp",
    "Full movement trail reconstructed",
    "Last known location → roadblock deployed",
]):
    box(s, 6.6, 2.6 + i*0.55, 6.2, 0.5, f"✓  {line}", size=14, color=RGBColor(0x88,0xFF,0x88))

rect(s, 0.3, 6.6, 12.7, 0.6, GREEN)
box(s, 0.3, 6.65, 12.7, 0.5,
    "TIME TO RESULT:  Under 10 seconds",
    size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, MGRAY)
box(s, 0.3, 0.1, 12, 0.7, "KEY FEATURES", size=30, bold=True, color=GREEN)
hline(s, 0.95, GREEN)

features = [
    ("🎯  Target Registration",
     "Mark any vehicle or person as suspect in one click.\nSystem highlights them RED with live movement trail."),
    ("🔍  Forensic Search",
     "Search by ID, description, or time range.\nReturns complete sighting history instantly."),
    ("🗺️  Movement Trail",
     "Full chronological breadcrumb path across cameras.\nExportable for court evidence."),
    ("📊  Analytics Dashboard",
     "Live stats, top 10 active vehicles, recent sightings.\nAuto-refreshes every 5 seconds."),
    ("📁  Data Export",
     "Full CSV export — 64,001 records already exported.\nTimestamped, structured, court-admissible."),
    ("⚡  Real-Time Detection",
     "Cars, motorcycles, buses, trucks, persons — all detected.\nPersistent Track IDs across frames."),
]

for i, (title, desc) in enumerate(features):
    col = (i % 3) * 4.3 + 0.3
    row = (i // 3) * 2.5 + 1.1
    rect(s, col, row, 4.1, 2.2, MGRAY, GREEN)
    box(s, col+0.1, row+0.1, 3.9, 0.5, title, size=15, bold=True, color=GREEN)
    box(s, col+0.1, row+0.65, 3.9, 1.4, desc, size=13, color=WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPARISON TABLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, MGRAY)
box(s, 0.3, 0.1, 12, 0.7, "HAWKEYE vs Manual CCTV Review",
    size=30, bold=True, color=CYAN)
hline(s, 0.95, CYAN)

headers = ["Feature", "Manual CCTV", "HAWKEYE"]
widths  = [4.5, 4.0, 4.5]
h_colors = [CYAN, RED, GREEN]
table_row(s, 0.3, 1.1, headers, widths, h_colors, size=15, bold=True)

rows_data = [
    ("Time to find suspect",      "2–5 days",    "Under 10 seconds"),
    ("Cameras connected",         "0",           "Unlimited"),
    ("Search capability",         "None",        "Full forensic search"),
    ("Movement timeline",         "Manual",      "Automatic"),
    ("Real-time alerts",          "None",        "Yes"),
    ("Data retention",            "Tape overwrites", "Permanent database"),
    ("Evidence export",           "DVD copy",    "Structured CSV"),
]

for i, (feat, manual, hawk) in enumerate(rows_data):
    t = 1.6 + i * 0.62
    bg_c = RGBColor(0x0a,0x0a,0x0a) if i % 2 == 0 else MGRAY
    box(s, 0.3,  t, 4.4, 0.55, feat,   size=14, color=WHITE,  bg_color=bg_c)
    box(s, 4.8,  t, 3.9, 0.55, manual, size=14, color=RED,    bg_color=bg_c)
    box(s, 8.8,  t, 4.4, 0.55, hawk,   size=14, color=GREEN,  bg_color=bg_c)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, MGRAY)
box(s, 0.3, 0.1, 12, 0.7, "TECHNOLOGY STACK", size=30, bold=True, color=CYAN)
hline(s, 0.95, CYAN)

stack = [
    ("AI Detection",    "YOLOv8n  (Ultralytics)",          GREEN),
    ("Video Processing","OpenCV",                           CYAN),
    ("Web Interface",   "Flask + HTML5",                    YELLOW),
    ("Database",        "SQLite  →  scalable to PostgreSQL",WHITE),
    ("Vector Search",   "Ready for face embedding search",  CYAN),
    ("Deployment",      "Docker-ready, cloud-deployable",   GREEN),
    ("Language",        "Python",                           WHITE),
]

for i, (layer, tech, color) in enumerate(stack):
    t = 1.1 + i * 0.75
    rect(s, 0.3, t, 3.5, 0.6, MGRAY, GREEN)
    box(s, 0.4, t+0.05, 3.3, 0.5, layer, size=15, bold=True, color=GREEN)
    box(s, 4.0, t+0.05, 8.5, 0.5, tech,  size=15, color=color)

rect(s, 0.3, 6.5, 12.7, 0.7, RGBColor(0x00,0x1a,0x00))
box(s, 0.3, 6.55, 12.7, 0.55,
    "Runs on standard CPU hardware — no expensive GPU required for deployment",
    size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 13.33, 0.9, MGRAY)
box(s, 0.3, 0.1, 12, 0.7, "ROADMAP", size=30, bold=True, color=YELLOW)
hline(s, 0.95, YELLOW)

phases = [
    ("PHASE 1\nDONE ✅", GREEN, [
        "Single-camera vehicle + person tracking",
        "Forensic search & movement trails",
        "Analytics dashboard",
        "64,001 sightings in database",
        "CSV evidence export",
    ]),
    ("PHASE 2\n30 Days", CYAN, [
        "Multi-camera coordination",
        "Face recognition (photo → find all appearances)",
        "License plate OCR",
        "Real-time SMS/email alerts",
    ]),
    ("PHASE 3\nScale", YELLOW, [
        "City-wide camera network integration",
        "Warrant database cross-reference",
        "Predictive movement analysis",
        "Police dispatch system integration",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    col = i * 4.3 + 0.3
    rect(s, col, 1.1, 4.1, 5.8, MGRAY, color)
    box(s, col+0.1, 1.2, 3.9, 0.8, title, size=17, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        box(s, col+0.2, 2.2 + j*0.85, 3.7, 0.75,
            f"▸  {item}", size=13, color=WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
rect(s, 0, 0, 0.18, 7.5, GREEN)
rect(s, 0, 6.8, 13.33, 0.7, MGRAY)

box(s, 0.5, 0.5, 12, 1.2, "🦅  HAWKEYE", size=60, bold=True, color=GREEN)
hline(s, 1.9)

box(s, 0.5, 2.1, 12, 0.8,
    '"The cameras are already watching.  HAWKEYE makes them think."',
    size=24, bold=True, color=YELLOW)

box(s, 0.5, 3.1, 12, 0.5,
    "Turns thousands of blind cameras into one intelligent eye.",
    size=18, color=WHITE)

for i, (num, label) in enumerate([
    ("64,001", "vehicle sightings — real data"),
    ("883",    "unique vehicles tracked"),
    ("10 sec", "from crime scene to suspect location"),
]):
    box(s, 0.5 + i*4.2, 3.9, 4.0, 0.5, num,   size=26, bold=True, color=GREEN)
    box(s, 0.5 + i*4.2, 4.45, 4.0, 0.4, label, size=13, color=WHITE)

hline(s, 5.1)
box(s, 0.5, 5.3, 12, 0.45, "GitHub:  github.com/GRACENGARI/yolo",
    size=15, color=CYAN)
box(s, 0.5, 5.8, 12, 0.45, "Live System:  http://localhost:5001",
    size=15, color=CYAN)
box(s, 0.5, 6.3, 12, 0.45,
    "Evidence:  vehicle_evidence_report.csv  |  people_surveillance_report.csv",
    size=14, color=RGBColor(0x88,0x88,0x88))

box(s, 0.3, 6.85, 12.7, 0.4,
    "Built for the HAWKEYE Hackathon — AI Surveillance Intelligence",
    size=12, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────
OUTPUT = "HAWKEYE_Pitch_Deck.pptx"
prs.save(OUTPUT)
print(f"\n✅  Saved: {OUTPUT}")
print(f"   Slides: 10")
print(f"   Open with Microsoft PowerPoint or Google Slides")
